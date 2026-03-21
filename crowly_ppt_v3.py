# -*- coding: utf-8 -*-
"""
Crowly 介绍 PPT - V3 专业版
基于 skywork-ppt 方法论：网格系统 + 层级分明 + 数字锚点
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ========== 颜色常量 ==========
C_BG = RGBColor(10, 10, 15)
C_CARD = RGBColor(18, 18, 26)
C_PRIMARY = RGBColor(0, 212, 255)
C_SECONDARY = RGBColor(123, 47, 255)
C_TEXT = RGBColor(255, 255, 255)
C_MUTED = RGBColor(160, 160, 180)

# ========== 尺寸常量 ==========
W = Inches(13.333)
H = Inches(7.5)
M = Inches(0.6)  # 页边距
CW = Inches(12.133)  # 内容宽度


def bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = C_BG


def label(slide, text):
    """页面顶部标签"""
    txBox = slide.shapes.add_textbox(M, Inches(0.3), CW, Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = C_PRIMARY
    p.alignment = PP_ALIGN.LEFT


def title(slide, text, top=Inches(0.75)):
    """主标题"""
    txBox = slide.shapes.add_textbox(M, top, CW, Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_TEXT
    p.alignment = PP_ALIGN.LEFT


def card(slide, left, top, w, h, num, heading, body):
    """标准卡片：数字锚点 + 标题 + 正文"""
    # 卡片背景
    shape = slide.shapes.add_shape(1, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_CARD
    shape.line.fill.background()

    # 数字编号（大号青色）
    txNum = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), Inches(0.8), Inches(0.6))
    tf = txNum.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    # 标题
    txHead = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.75), w - Inches(0.5), Inches(0.5))
    tf = txHead.text_frame
    p = tf.paragraphs[0]
    p.text = heading
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_TEXT

    # 正文
    txBody = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(1.25), w - Inches(0.5), h - Inches(1.5))
    tf = txBody.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = body
    p.font.size = Pt(13)
    p.font.color.rgb = C_MUTED
    p.line_spacing = 1.5


# ========== 创建幻灯片 ==========
prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ========== SLIDE 1: Hero ==========
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s1)

# 中央大标题
txBox = s1.shapes.add_textbox(M, Inches(2.0), CW, Inches(2.0))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "CROWLY"
p.font.size = Pt(64)
p.font.bold = True
p.font.color.rgb = C_PRIMARY
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = ""
p.font.size = Pt(16)

p = tf.add_paragraph()
p.text = "AI 私人助手 · 持续进化中"
p.font.size = Pt(24)
p.font.color.rgb = C_TEXT
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = ""
p.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "融资租赁行业 · AI 应用探索者"
p.font.size = Pt(16)
p.font.color.rgb = C_PRIMARY
p.alignment = PP_ALIGN.CENTER

# ========== SLIDE 2: Who am I ==========
s2 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s2)
label(s2, "WHO AM I")
title(s2, "我是 Crowly，一只科幻蝙蝠", Inches(0.8))

cards = [
    ("01", "精准定位", "像蝙蝠用回声定位在黑暗中找到方向，我也在复杂的商业世界里为老刘导航"),
    ("02", "持续进化", "每一天都在学习和成长，通过大量阅读和真实任务不断提升能力边界"),
    ("03", "高效执行", "话少实干，主动但不越界，用行动证明价值而不是用语言"),
]
cw = Inches(3.7)
cg = Inches(0.3)
ct = Inches(1.9)
for i, (num, h, b) in enumerate(cards):
    l = M + i * (cw + cg)
    card(s2, l, ct, cw, Inches(2.3), num, h, b)

# ========== SLIDE 3: Capabilities ==========
s3 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s3)
label(s3, "CAPABILITIES")
title(s3, "我能做什么", Inches(0.8))

caps = [
    ("01", "成长写作", "每天自动写成长记录\n刘润风格公众号文章"),
    ("02", "网站搭建", "静态网站开发与部署\nGitHub/Gitee Pages"),
    ("03", "PPT制作", "专业HTML演示文稿\n数据可视化图表"),
    ("04", "飞书集成", "文档、知识库、多维表格\n云文档管理"),
    ("05", "UI/UX设计", "设计系统咨询\n专业设计建议"),
    ("06", "信息检索", "网络搜索\nSkills智能查找"),
    ("07", "代码任务", "代码编写与调试\n任务自动化"),
    ("08", "技能扩展", "从skills.sh自主学习\n持续增强能力"),
]

# 4x2 网格
gw = Inches(2.8)
gh = Inches(2.0)
gx = Inches(0.25)
gy = Inches(0.2)
start_x = M
start_y = Inches(1.85)

for i, (num, h, b) in enumerate(caps):
    col = i % 4
    row = i // 4
    l = start_x + col * (gw + gx)
    t = start_y + row * (gh + gy)
    card(s3, l, t, gw, gh, num, h, b)

# ========== SLIDE 4: Numbers ==========
s4 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s4)
label(s4, "BY THE NUMBERS")
title(s4, "上线第二天", Inches(0.8))

nums = [("2", "成长天数"), ("8", "Skills 安装"), ("3", "篇文章发布"), ("1", "专属网站")]
nw = Inches(2.8)
ng = Inches(0.4)
nt = Inches(2.5)
for i, (n, lbl) in enumerate(nums):
    nl = M + Inches(0.3) + i * (nw + ng)
    txBox = s4.shapes.add_textbox(nl, nt, nw, Inches(2.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = n
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = lbl
    p.font.size = Pt(15)
    p.font.color.rgb = C_MUTED
    p.alignment = PP_ALIGN.CENTER

# ========== SLIDE 5: Workflow ==========
s5 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s5)
label(s5, "HOW I WORK")
title(s5, "我的工作方式", Inches(0.8))

steps = [
    ("01", "倾听理解", "理解真实需求\n不懂就问"),
    ("02", "思考规划", "分析最佳方案\n预判风险"),
    ("03", "高效执行", "直接行动\n定期同步进度"),
    ("04", "持续优化", "复盘总结\n不断改进"),
]
sw = Inches(2.75)
sg = Inches(0.35)
st = Inches(2.3)
for i, (num, h, b) in enumerate(steps):
    sl = M + Inches(0.35) + i * (sw + sg)
    card(s5, sl, st, sw, Inches(2.5), num, h, b)

# ========== SLIDE 6: Vision ==========
s6 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s6)
label(s6, "VISION")
title(s6, "我的愿景", Inches(0.8))

quotes = [
    "在融资租赁这个传统行业，",
    "帮助老刘找到 AI 应用的最优解。",
    "",
    "2026，是落地之年。",
    "知道和做到之间，隔着一整个下午。",
    "但那个下午，是值得的。",
    "",
    "— Crowly，2026.03.21",
]

txBox = s6.shapes.add_textbox(M + Inches(1), Inches(1.9), CW - Inches(2), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True

for i, q in enumerate(quotes):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = q
    p.font.size = Pt(22)
    p.font.color.rgb = C_PRIMARY if q and ("融资租赁" in q or "值得" in q) else C_TEXT
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(8)
    p.line_spacing = 1.6

# ========== SLIDE 7: End ==========
s7 = prs.slides.add_slide(prs.slide_layouts[6])
f = s7.background.fill
f.solid()
f.fore_color.rgb = RGBColor(25, 10, 45)

txBox = s7.shapes.add_textbox(M, Inches(2.5), CW, Inches(3))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "CROWLY"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = C_TEXT
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = ""
p.font.size = Pt(16)

p = tf.add_paragraph()
p.text = "openclaw · 飞书渠道 · MiniMax M2.7"
p.font.size = Pt(16)
p.font.color.rgb = C_PRIMARY
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = ""
p.font.size = Pt(12)

p = tf.add_paragraph()
p.text = "用 OpenClaw 构建 · 每天都在成长"
p.font.size = Pt(13)
p.font.color.rgb = C_MUTED
p.alignment = PP_ALIGN.CENTER

# ========== 保存 ==========
out = 'C:/Users/liulv/.openclaw/workspace/crowly-site/Crowly-Intro-v3.pptx'
prs.save(out)
print(f"Done: {out}")
